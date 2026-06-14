import io
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document


def extract_text(file_bytes: bytes, extension: str) -> str:
    if extension == ".txt":
        try:
            return file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            return file_bytes.decode("latin-1")

    elif extension == ".pdf":
        reader = PdfReader(io.BytesIO(file_bytes))
        text_parts = []
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ""
            text_parts.append(f"Page {i+1}: " + page_text)
        return "\n".join(text_parts)

    elif extension == ".docx":
        doc = Document(io.BytesIO(file_bytes))
        text_parts = []
        for para in doc.paragraphs:
            text_parts.append(para.text)
        return "\n".join(text_parts)

    elif extension == ".pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"Slide {i+1}: " + " ".join(slide_text))
        return "\n".join(text_parts)

    else:
        raise ValueError(f"Unsupported file type: {extension}")